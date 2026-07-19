from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(r"D:\BTL_HTN_v2")
PROJECT = ROOT / "projects" / "BTL_slide_topic1_ppt169_20260718"
TEMPLATE = PROJECT / "sources" / "BTL_slide.pptx"
OUTPUT = PROJECT / "exports" / "BTL_slide_hoan_chinh.pptx"
FIGS = ROOT / "report" / "generated" / "la_testsuite_20260718" / "figs"

NAVY = RGBColor(0x00, 0x3B, 0x70)
NAVY_DARK = RGBColor(0x00, 0x2D, 0x57)
RED = RGBColor(0xC9, 0x25, 0x35)
BLUE = RGBColor(0x1C, 0x6E, 0xA4)
CYAN = RGBColor(0x18, 0x92, 0xB8)
GREEN = RGBColor(0x1B, 0x8A, 0x5A)
AMBER = RGBColor(0xD2, 0x78, 0x14)
INK = RGBColor(0x1E, 0x2A, 0x35)
MUTED = RGBColor(0x5D, 0x6B, 0x78)
LIGHT = RGBColor(0xF1, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def clone_template_slide(prs: Presentation, template_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for rel in template_slide.part.rels.values():
        if rel.reltype.endswith("/image") and not rel.is_external:
            slide.part.rels._add_relationship(rel.reltype, rel.target_part)
    for shape in template_slide.shapes:
        slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return slide


def add_fade(slide):
    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))
    slide._element.insert(2, transition)


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Arial", valign=MSO_ANCHOR.TOP,
             margin=0.05, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return box


def add_title(slide, title, number):
    add_text(slide, title.upper(), 0.75, 0.27, 16.8, 0.55, 27, WHITE, True,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, f"{number:02d}", 18.1, 0.27, 1.0, 0.5, 18, WHITE, True,
             PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def add_footer_label(slide, text="THIẾT BỊ LOGIC ANALYZER 8 KÊNH"):
    add_text(slide, text, 13.7, 10.50, 5.1, 0.25, 9, NAVY, True, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body="", accent=NAVY, title_size=18,
             body_size=15, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xD4, 0xDF, 0xE7)
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.11), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    add_text(slide, title, x + 0.24, y + 0.18, w - 0.42, 0.42, title_size, accent, True)
    if body:
        add_text(slide, body, x + 0.24, y + 0.68, w - 0.42, h - 0.82, body_size, INK)
    return shape


def add_kpi(slide, x, y, w, h, value, label, accent=NAVY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = RGBColor(0xD4, 0xDF, 0xE7)
    add_text(slide, value, x + 0.15, y + 0.16, w - 0.3, 0.62, 27, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.18, y + 0.82, w - 0.36, h - 0.95, 13, MUTED, False, PP_ALIGN.CENTER)


def add_bullets(slide, items, x, y, w, h, size=18, color=INK, gap=0.58, accent=RED):
    for i, item in enumerate(items):
        yy = y + i * gap
        add_text(slide, "•", x, yy, 0.26, 0.3, size + 3, accent, True)
        add_text(slide, item, x + 0.32, yy, w - 0.32, min(gap, h - i * gap), size, color)


def add_arrow(slide, x, y, w, h, color=NAVY):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()


def add_picture_fit(slide, path, x, y, w, h, border=True):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    target = w / h
    source = iw / ih
    if source > target:
        pic_h = h
        pic_w = h * source
        left = x - (pic_w - w) / 2
        top = y
    else:
        pic_w = w
        pic_h = w / source
        left = x
        top = y - (pic_h - h) / 2
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(pic_w), Inches(pic_h))
    if source > target:
        crop = (pic_w - w) / pic_w / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (pic_h - h) / pic_h / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    if border:
        pic.line.color.rgb = RGBColor(0xC4, 0xD1, 0xDC)
        pic.line.width = Pt(1)
    return pic


def add_picture_contain(slide, path, x, y, w, h, border=True):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    source = iw / ih
    target = w / h
    if source > target:
        pic_w = w
        pic_h = w / source
        left = x
        top = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * source
        left = x + (w - pic_w) / 2
        top = y
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(pic_w), Inches(pic_h))
    if border:
        pic.line.color.rgb = RGBColor(0xC4, 0xD1, 0xDC)
        pic.line.width = Pt(1)
    return pic


def add_caption(slide, text, x, y, w):
    add_text(slide, text, x, y, w, 0.32, 11, MUTED, False, PP_ALIGN.CENTER, italic=True)


def set_notes(slide, text):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text


def build():
    prs = Presentation(str(TEMPLATE))
    cover = prs.slides[0]
    base = prs.slides[1]
    while len(prs.slides) < 15:
        clone_template_slide(prs, base)
    slides = list(prs.slides)
    for slide in slides:
        add_fade(slide)

    # 1. Cover
    add_text(cover, "BÁO CÁO BÀI TẬP LỚN", 1.05, 2.38, 10.8, 0.55, 22, RED, True)
    add_text(cover, "THIẾT KẾ VÀ XÂY DỰNG\nTHIẾT BỊ LOGIC ANALYZER ĐƠN GIẢN",
             1.05, 3.05, 11.8, 1.65, 31, NAVY_DARK, True)
    add_text(cover, "Đề tài 1 · Hệ thống nhúng và thiết kế giao tiếp nhúng",
             1.05, 4.82, 11.5, 0.48, 17, MUTED)
    add_text(cover, "Đoàn Sinh Đức — 20234000\nPhạm Đăng Vinh — 20233719\nVũ Mạnh Quân — 20234033\nVũ Nam Khánh — 20234015",
             1.05, 6.00, 6.1, 1.55, 16, INK)
    add_text(cover, "GIẢNG VIÊN HƯỚNG DẪN\nTS. ĐÀO VIỆT HÙNG", 7.35, 6.00, 4.7, 1.0, 15, NAVY, True)
    add_text(cover, "Hà Nội, 7/2026", 1.05, 8.15, 3.5, 0.35, 13, MUTED)
    set_notes(cover, "Nhóm trình bày kết quả thiết kế và xây dựng thiết bị phân tích logic tám kênh trên STM32F103C8. Nội dung tập trung vào kiến trúc, cơ chế lấy mẫu, phần mềm hiển thị và các kết quả kiểm thử phần cứng trong vòng.")

    # 2. Problem
    s = slides[1]; add_title(s, "Đặt vấn đề", 2); add_footer_label(s)
    add_text(s, "LỖI SỐ KHÔNG CHỈ NẰM Ở MỨC 0/1", 1.05, 1.65, 8.6, 0.62, 27, NAVY_DARK, True)
    add_bullets(s, [
        "Quan hệ thời gian giữa các kênh quyết định tính đúng của giao tiếp.",
        "Thiếu xung, sai thứ tự cạnh hoặc sai khung rất khó quan sát bằng đồng hồ đo.",
        "Cần ghi đồng thời nhiều tín hiệu và xem lại theo trục thời gian trên máy tính."
    ], 1.05, 2.65, 9.2, 3.0, 19, gap=0.9)
    add_card(s, 11.1, 1.85, 3.7, 2.15, "GHI NHẬN", "8 đường tín hiệu\ntại cùng thời điểm", BLUE, 18, 17)
    add_arrow(s, 15.0, 2.55, 1.0, 0.45, RED)
    add_card(s, 16.15, 1.85, 2.8, 2.15, "PHÂN TÍCH", "Waveform\n+ decoder", GREEN, 18, 17)
    add_text(s, "Logic analyzer biến chuỗi trạng thái rời rạc thành bằng chứng về timing và giao thức.", 11.1, 4.55, 7.85, 1.2, 21, NAVY, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    set_notes(s, "Trong hệ thống số, mức điện áp đúng chưa đủ để kết luận giao tiếp đúng. Nhiều lỗi chỉ xuất hiện trong quan hệ thời gian giữa các đường tín hiệu, vì vậy nhóm cần một công cụ ghi đồng thời nhiều kênh và cho phép phân tích lại trên máy tính.")

    # 3. Goals
    s = slides[2]; add_title(s, "Mục tiêu và phạm vi", 3); add_footer_label(s)
    goals = [
        ("8 KÊNH", "CH0–CH7\nPA0–PA7", BLUE),
        ("TỪ 1 kHz", "Cấu hình tốc độ\nlấy mẫu", GREEN),
        ("THU NGOẠI TUYẾN", "Buffer → DUMP\nkhung SLA8", AMBER),
        ("PC GUI", "Waveform, đo cạnh,\nzoom và cuộn", CYAN),
        ("3 DECODER", "UART · I2C · SPI", RED),
    ]
    for i, (a, b, c) in enumerate(goals):
        x = 0.9 + i * 3.75
        add_card(s, x, 1.85, 3.35, 2.15, a, b, c, 19, 16)
    add_text(s, "PHẠM VI TRIỂN KHAI", 1.0, 4.75, 4.2, 0.45, 20, NAVY, True)
    add_bullets(s, [
        "Firmware STM32: TIM2, DMA/ISR, trigger và giao thức lệnh.",
        "Phần mềm PC: kiểm khung, hiển thị tám kênh và giải mã giao thức.",
        "Arduino UNO tạo tín hiệu Gray, UART, I2C và SPI để kiểm thử từng chế độ."
    ], 1.0, 5.35, 17.5, 2.5, 18, gap=0.78)
    set_notes(s, "Mục tiêu tối thiểu của đề tài là hai kênh và tốc độ từ một kilohertz. Mẫu thử thực tế được triển khai với tám kênh, có giao diện máy tính và ba bộ giải mã giao thức. Phạm vi đánh giá gồm cả firmware, khung truyền dữ liệu và phần mềm PC.")

    # 4. Architecture
    s = slides[3]; add_title(s, "Kiến trúc tổng thể", 4); add_footer_label(s)
    add_card(s, 0.9, 2.0, 4.2, 4.3, "NGUỒN TÍN HIỆU", "Arduino UNO\nGray · UART · I2C · SPI\n\nOpen-drain, chung GND", AMBER, 20, 18)
    add_arrow(s, 5.25, 3.75, 1.0, 0.5, RED)
    add_card(s, 6.35, 1.65, 6.2, 5.0, "STM32F103C8", "PA0–PA7 → GPIOA IDR\nTIM2 tạo nhịp lấy mẫu\nDMA hoặc ISR → buffer 13.888 mẫu\nSLA8: metadata + payload + checksum", NAVY, 22, 18)
    add_arrow(s, 12.75, 3.75, 1.0, 0.5, RED)
    add_card(s, 13.85, 2.0, 5.2, 4.3, "MÁY TÍNH", "USART1 · 1 Mbaud\nPyQt5 + pyqtgraph\nWaveform 8 kênh\nUART · I2C · SPI decoder", GREEN, 20, 18)
    add_text(s, "Một mẫu = một lần đọc đồng thời 8 bit thấp của GPIOA", 4.7, 7.25, 10.7, 0.65, 20, NAVY, True, PP_ALIGN.CENTER)
    set_notes(s, "Kiến trúc được chia thành ba lớp. Arduino tạo tín hiệu tham chiếu, STM32 lấy mẫu tám bit đồng thời và đóng gói thành khung SLA8, còn phần mềm PC kiểm tra khung rồi hiển thị hoặc giải mã. Việc đọc GPIOA IDR một lần giúp giảm sai lệch tương đối giữa các kênh.")

    # 5. Hardware
    s = slides[4]; add_title(s, "Phần cứng và đấu nối", 5); add_footer_label(s)
    add_card(s, 0.9, 1.75, 5.7, 5.7, "STM32F103C8", "• PA0–PA7: CH0–CH7\n• PA9/PA10: USART1\n• Timer clock: 72 MHz\n• Logic: 3,3 V\n• Buffer: 13.888 byte", NAVY, 22, 18)
    add_card(s, 7.15, 1.75, 5.7, 5.7, "ARDUINO UNO", "• D2–D9 → PA0–PA7\n• Tạo mẫu Gray 8 bit\n• Phát UART/I2C/SPI riêng\n• Ngõ ra open-drain\n• Dùng pull-up 3,3 V", AMBER, 22, 18)
    add_card(s, 13.4, 1.75, 5.7, 5.7, "NGUYÊN TẮC AN TOÀN", "• Nối chung GND\n• Không đưa 5 V trực tiếp\n• HIGH nhờ pull-up 3,3 V\n• Kiểm tra mức trước khi test\n• Không suy diễn mạch bảo vệ", RED, 22, 18)
    add_text(s, "Không khẳng định có khối bảo vệ đầu vào khi chưa có sơ đồ nguyên lý được kiểm chứng.", 2.1, 8.05, 15.8, 0.55, 18, RED, True, PP_ALIGN.CENTER)
    set_notes(s, "Tám đầu vào được đưa trực tiếp vào PA0 đến PA7, còn USART1 dùng để truyền dữ liệu về máy tính. Arduino được cấu hình open-drain để mức cao do điện trở kéo lên ba phẩy ba volt tạo ra. Báo cáo không giả định thêm khối bảo vệ đầu vào vì chưa có sơ đồ nguyên lý đầy đủ.")

    # 6. Sampling
    s = slides[5]; add_title(s, "Cơ chế lấy mẫu DMA và ISR", 6); add_footer_label(s)
    add_text(s, "TRIGGER TỨC THỜI", 1.0, 1.65, 8.5, 0.45, 21, NAVY, True)
    add_card(s, 1.0, 2.25, 3.1, 1.45, "TIM2 UPDATE", "Nhịp lấy mẫu", BLUE, 18, 15)
    add_arrow(s, 4.25, 2.75, 0.8, 0.4)
    add_card(s, 5.15, 2.25, 3.1, 1.45, "DMA1 CH2", "GPIOA IDR → RAM", GREEN, 18, 15)
    add_arrow(s, 8.4, 2.75, 0.8, 0.4)
    add_card(s, 9.3, 2.25, 3.1, 1.45, "BUFFER ĐỦ", "Phát EVENT", RED, 18, 15)
    add_text(s, "TRIGGER CẠNH / MẪU", 1.0, 4.65, 8.5, 0.45, 21, NAVY, True)
    add_card(s, 1.0, 5.25, 3.1, 1.45, "TIM2 ISR", "Đọc và kiểm điều kiện", AMBER, 18, 15)
    add_arrow(s, 4.25, 5.75, 0.8, 0.4)
    add_card(s, 5.15, 5.25, 3.1, 1.45, "PRE-TRIGGER", "Bộ đệm vòng", CYAN, 18, 15)
    add_arrow(s, 8.4, 5.75, 0.8, 0.4)
    add_card(s, 9.3, 5.25, 3.1, 1.45, "POST-TRIGGER", "Hoàn thành capture", RED, 18, 15)
    add_kpi(s, 14.0, 2.15, 4.4, 1.75, "6,545 MS/s", "Trần DMA được kiểm chứng", NAVY)
    add_kpi(s, 14.0, 4.55, 4.4, 1.75, "400 kS/s", "Trần ISR được kiểm chứng", AMBER)
    set_notes(s, "Với trigger tức thời, DMA chuyển trực tiếp trạng thái GPIO vào RAM nên đạt tốc độ cao hơn. Khi cần trigger cạnh hoặc mẫu, ISR phải kiểm tra điều kiện và quản lý vùng trước sau trigger, vì vậy trần tốc độ thấp hơn. Hai giới hạn trên đều được xác nhận bằng kiểm thử phần cứng.")

    # 7. Protocol
    s = slides[6]; add_title(s, "Firmware và khung dữ liệu SLA8", 7); add_footer_label(s)
    flow = [("CFG", BLUE), ("ARM", NAVY), ("CAPTURE", AMBER), ("EVENT", RED), ("DUMP", GREEN)]
    for i, (label, col) in enumerate(flow):
        x = 0.9 + i * 3.75
        add_card(s, x, 1.65, 3.0, 1.25, label, "", col, 20)
        if i < len(flow)-1: add_arrow(s, x+3.1, 2.05, 0.5, 0.3, RED)
    add_text(s, "KHUNG SLA8", 0.95, 3.65, 3.0, 0.45, 22, NAVY, True)
    widths = [4.0, 2.6, 7.5, 3.2]
    labels = [("HEADER 48 B", NAVY), ("METADATA", BLUE), ("PAYLOAD ≤ 13.888 B", GREEN), ("CHECKSUM", RED)]
    x = 0.95
    for w, (lab, col) in zip(widths, labels):
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.35), Inches(w), Inches(1.25))
        box.fill.solid(); box.fill.fore_color.rgb = col; box.line.fill.background()
        add_text(s, lab, x+0.08, 4.68, w-0.16, 0.42, 17, WHITE, True, PP_ALIGN.CENTER)
        x += w
    add_bullets(s, [
        "1 byte cho mỗi mẫu tám kênh.",
        "Checksum FNV-1a cho header và payload.",
        "Metadata lưu rate, mode, trigger và trạng thái lỗi."
    ], 1.0, 6.55, 17.0, 1.8, 17, gap=0.55)
    set_notes(s, "Firmware làm việc theo chuỗi cấu hình, arm, capture, báo sự kiện rồi mới dump dữ liệu. Khung SLA8 tách metadata khỏi payload và dùng checksum FNV một a để phát hiện lỗi truyền. Mỗi mẫu chỉ chiếm một byte vì tám kênh được đóng vào tám bit.")

    # 8. PC GUI
    s = slides[7]; add_title(s, "Phần mềm PC và giao diện", 8); add_footer_label(s)
    add_picture_contain(s, FIGS / "fig_gui_overview.png", 0.85, 1.45, 12.7, 7.65)
    add_card(s, 14.0, 1.6, 5.0, 1.65, "KẾT NỐI", "INFO · STATUS · cấu hình", BLUE, 18, 15)
    add_card(s, 14.0, 3.55, 5.0, 1.65, "WAVEFORM", "8 kênh · zoom · cuộn", GREEN, 18, 15)
    add_card(s, 14.0, 5.50, 5.0, 1.65, "PHÂN TÍCH", "Đo cạnh · tần số · decoder", RED, 18, 15)
    add_caption(s, "Trạng thái trong ảnh xác nhận SLA8-FW trên COM12; dropdown Port chỉ là lựa chọn giao diện", 1.0, 9.22, 12.4)
    set_notes(s, "Phần mềm PC chịu trách nhiệm cấu hình thiết bị, nhận khung và kiểm checksum trước khi hiển thị. Người dùng có thể xem tám kênh, phóng to theo thời gian, đo cạnh và chạy các decoder UART, I2C hoặc SPI trên dữ liệu đã thu.")

    # 9. HIL method
    s = slides[8]; add_title(s, "Phương pháp kiểm thử phần cứng trong vòng", 9); add_footer_label(s)
    add_card(s, 0.9, 1.6, 4.0, 1.5, "ARDUINO · COM18", "Tín hiệu tham chiếu", AMBER, 20, 16)
    add_arrow(s, 5.0, 2.12, 0.8, 0.4, RED)
    add_card(s, 5.9, 1.6, 4.0, 1.5, "STM32 · COM12", "Thu tám kênh", NAVY, 20, 16)
    add_arrow(s, 10.0, 2.12, 0.8, 0.4, RED)
    add_card(s, 10.9, 1.6, 4.0, 1.5, "SLA8 + CHECKSUM", "Khung bằng chứng", GREEN, 20, 16)
    add_arrow(s, 15.0, 2.12, 0.8, 0.4, RED)
    add_card(s, 15.9, 1.6, 3.2, 1.5, "ORACLE", "Đối chiếu tự động", RED, 20, 15)
    groups = [
        ("TOÀN VẸN", "TC01–TC02\nKhung, 8 kênh, Gray", BLUE),
        ("TỐC ĐỘ", "TC03–TC06\nRate, trần, Nyquist", NAVY),
        ("GIAO THỨC", "TC07–TC09\nUART, I2C, SPI", GREEN),
        ("TRIGGER", "TC10\nCạnh rơi và đầu vào sai", RED),
    ]
    for i,(a,b,c) in enumerate(groups):
        add_card(s, 0.9+i*4.65, 4.3, 4.1, 2.35, a, b, c, 20, 17)
    add_text(s, "Gray 8 bit: mỗi bước hợp lệ chỉ thay đổi đúng một bit", 3.2, 7.45, 13.5, 0.55, 20, NAVY, True, PP_ALIGN.CENTER)
    set_notes(s, "Kiểm thử phần cứng trong vòng dùng Arduino làm nguồn tham chiếu và STM32 làm thiết bị cần đánh giá. Chuỗi Gray phù hợp để phát hiện mất mẫu hoặc sai quan hệ giữa kênh vì hai trạng thái liên tiếp chỉ được phép khác một bit. Mười kịch bản được nhóm theo toàn vẹn, tốc độ, giao thức và trigger.")

    # 10. Speed
    s = slides[9]; add_title(s, "Kết quả tốc độ và toàn vẹn dữ liệu", 10); add_footer_label(s)
    add_picture_contain(s, FIGS / "fig_tc04_ceiling.png", 0.85, 1.45, 12.3, 7.35)
    add_kpi(s, 13.75, 1.55, 4.9, 1.75, "6,545 MS/s", "DMA sạch · sai số 0,0079%", NAVY)
    add_kpi(s, 13.75, 3.75, 4.9, 1.75, "400 kS/s", "ISR sạch · 0 overrun", AMBER)
    add_kpi(s, 13.75, 5.95, 4.9, 1.75, "0", "Overflow · dropped · checksum lỗi", GREEN)
    add_text(s, "7 / 8 / 10 MS/s bị firmware từ chối thay vì thu sai", 13.65, 8.15, 5.1, 0.65, 15, RED, True, PP_ALIGN.CENTER)
    add_caption(s, "TC-04 phát lại offline từ frame .sla8; dropdown Port không phải nguồn capture", 1.0, 9.05, 12.0)
    set_notes(s, "DMA thu sạch tại sáu phẩy năm bốn năm megasample mỗi giây với sai số dưới một phần trăm nghìn. ISR đạt bốn trăm kilosample mỗi giây và không ghi nhận overrun. Khi yêu cầu vượt trần, firmware trả lỗi thay vì tiếp tục với dữ liệu không đáng tin cậy.")

    # 11. Gray + aliasing
    s = slides[10]; add_title(s, "Tám kênh đồng thời và giới hạn Nyquist", 11); add_footer_label(s)
    add_picture_contain(s, FIGS / "fig_tc02_gray8ch.png", 0.75, 1.45, 8.95, 6.95)
    add_picture_contain(s, FIGS / "fig_tc06_aliasing.png", 10.25, 1.45, 8.95, 6.95)
    add_caption(s, "Gray 8 kênh: 0 lỗi chuỗi, 0 short run", 0.9, 8.55, 8.6)
    add_caption(s, "25 kHz: đúng ở 1 MS/s, alias khoảng 5 kHz ở 30 kS/s", 10.4, 8.55, 8.6)
    add_text(s, "Skew < 1 chu kỳ mẫu là chặn trên từ oracle Gray — không phải phép đo jitter tuyệt đối.", 2.1, 9.25, 15.8, 0.48, 16, RED, True, PP_ALIGN.CENTER)
    set_notes(s, "Dữ liệu Gray xác nhận đủ tám kênh và không phát hiện lỗi trình tự trong cửa sổ kiểm thử. Từ oracle này chỉ có thể kết luận sai lệch giữa kênh nhỏ hơn một chu kỳ mẫu, chưa thể coi là phép đo jitter tuyệt đối. Thử nghiệm Nyquist cho thấy lấy mẫu quá thấp biến tín hiệu hai mươi lăm kilohertz thành thành phần giả khoảng năm kilohertz.")

    # 12. Decoders
    s = slides[11]; add_title(s, "Giải mã UART, I2C và SPI", 12); add_footer_label(s)
    imgs = [
        ("UART · 57.600 baud", "55 A5 4F 4B · 0 framing error", FIGS/"fig_tc07_uart.png", BLUE),
        ("I2C · 2 MS/s", "START · 0x50 W · A5 ACK · 5A NACK · STOP", FIGS/"fig_tc08_i2c.png", GREEN),
        ("SPI · 500 kS/s", "MOSI/MISO: 55/A5 · A5/3C · 5A/C3", FIGS/"fig_tc09_spi.png", RED),
    ]
    for i,(title,body,img,col) in enumerate(imgs):
        x=0.65+i*6.45
        add_text(s,title,x,1.45,6.0,0.42,19,col,True,PP_ALIGN.CENTER)
        add_picture_contain(s,img,x,2.05,6.0,5.3)
        add_text(s,body,x,7.65,6.0,0.95,14,INK,True,PP_ALIGN.CENTER)
    add_text(s, "Ảnh phát lại offline từ .sla8; COM12/COM18 được xác minh riêng bằng device probe.", 3.0, 9.10, 14.0, 0.45, 16, NAVY, True, PP_ALIGN.CENTER)
    set_notes(s, "Ba bộ giải mã được kiểm tra bằng tín hiệu thật. UART thu đúng chuỗi byte và không có lỗi framing, I2C nhận đúng địa chỉ cùng trạng thái ACK và NACK, còn SPI nhận đúng ba cặp MOSI và MISO trong giao dịch có chip select.")

    # 13. Edge cases
    s = slides[12]; add_title(s, "Trigger và xử lý trường hợp biên", 13); add_footer_label(s)
    add_text(s, "SPI LẤY MẪU KHÔNG ĐỦ", 0.9, 1.45, 8.8, 0.42, 20, RED, True, PP_ALIGN.CENTER)
    add_picture_contain(s, FIGS / "fig_tc09_spi_undersampled.png", 0.8, 2.05, 9.0, 5.85)
    add_text(s, "Cảnh báo UNDERSAMPLED · không phát byte sai", 1.0, 8.15, 8.6, 0.4, 16, RED, True, PP_ALIGN.CENTER)
    add_text(s, "TRIGGER FALL CH6", 10.25, 1.45, 8.8, 0.42, 20, NAVY, True, PP_ALIGN.CENTER)
    add_picture_contain(s, FIGS / "fig_tc10_trigger.png", 10.2, 2.05, 9.0, 5.85)
    add_text(s, "Trigger index 1490 · vùng pre-trigger giữ mức HIGH", 10.4, 8.15, 8.6, 0.4, 16, NAVY, True, PP_ALIGN.CENTER)
    add_text(s, "Ảnh phát lại offline từ .sla8; đầu vào trigger/pattern sai được từ chối.", 4.0, 9.15, 12.0, 0.45, 16, GREEN, True, PP_ALIGN.CENTER)
    set_notes(s, "Thiết bị cần hành xử an toàn cả khi điều kiện đo không phù hợp. Với SPI bị lấy mẫu thiếu, decoder đưa ra cảnh báo và không tạo byte sai. Trigger cạnh rơi trên kênh sáu giữ được vùng trước trigger, còn cấu hình sai bị firmware từ chối rõ ràng.")

    # 14. Requirements
    s = slides[13]; add_title(s, "Đối chiếu yêu cầu và kết quả", 14); add_footer_label(s)
    cards = [
        ("8 KÊNH", "Yêu cầu ≥ 2", NAVY),
        ("1 kHz → 6,545 MS/s", "Yêu cầu từ 1 kHz", BLUE),
        ("GUI 8 KÊNH", "Waveform + đo cạnh", CYAN),
        ("UART · I2C · SPI", "Đã kiểm chứng vật lý", GREEN),
        ("10/10 TC", "Toàn bộ kịch bản đạt", RED),
        ("FAIL-SAFE", "Từ chối rate/cấu hình sai", AMBER),
    ]
    for i,(a,b,c) in enumerate(cards):
        row=i//3; col=i%3
        add_card(s, 0.9+col*6.25, 1.65+row*3.35, 5.65, 2.55, "✓  "+a, b, c, 21, 17)
    add_text(s, "Kết quả được chốt từ metrics.json và capture .sla8.", 3.0, 8.85, 14.0, 0.5, 17, NAVY, True, PP_ALIGN.CENTER)
    set_notes(s, "Mẫu thử vượt yêu cầu tối thiểu về số kênh và giữ được mức một kilohertz ở đầu dải. Giao diện hiển thị đủ tám kênh, ba decoder đều có bằng chứng vật lý và toàn bộ mười kịch bản kiểm thử đạt. Các trường hợp vượt giới hạn được từ chối thay vì âm thầm sinh dữ liệu sai.")

    # 15. Limitations/conclusion
    s = slides[14]; add_title(s, "Hạn chế, hướng phát triển và kết luận", 15); add_footer_label(s, "KẾT LUẬN")
    add_card(s, 0.9, 1.55, 8.7, 5.5, "HẠN CHẾ HIỆN TẠI", "• Chưa đo jitter tuyệt đối bằng chuẩn thời gian độc lập.\n• Chưa có sơ đồ nguyên lý đầu vào đầy đủ.\n• SPI mới hỗ trợ cạnh lên, MSB-first.\n• Trigger dùng ISR giới hạn ở 400 kS/s.", RED, 22, 18)
    add_card(s, 10.35, 1.55, 8.7, 5.5, "HƯỚNG PHÁT TRIỂN", "• Đo jitter/skew bằng tín hiệu tham chiếu chung.\n• Bổ sung CPOL, CPHA và thứ tự bit cho SPI.\n• Hoàn thiện sơ đồ đấu nối và quy trình kiểm mức.\n• Tối ưu trigger tốc độ cao.", GREEN, 22, 18)
    add_text(s, "Mẫu thử logic analyzer 8 kênh đã chạy HIL, giải mã ba giao thức và hành xử an toàn khi vượt giới hạn.", 1.6, 7.75, 16.8, 1.15, 23, NAVY_DARK, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "XIN CẢM ƠN!", 6.9, 9.25, 6.2, 0.55, 24, RED, True, PP_ALIGN.CENTER)
    set_notes(s, "Hệ thống vẫn còn các giới hạn cần được trình bày trung thực, đặc biệt là jitter tuyệt đối, mạch đầu vào và cấu hình SPI. Tuy vậy, mẫu thử đã hoàn thành chuỗi chức năng chính từ thu tám kênh, truyền khung, hiển thị đến giải mã và trigger. Kết quả HIL cho thấy hệ thống hoạt động ổn định trong miền đã kiểm chứng và từ chối các cấu hình vượt giới hạn.")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
